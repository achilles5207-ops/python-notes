import os
import docx
from pptx import Presentation

folder = "/workspaces/python-notes/resources"
for file in sorted(os.listdir(folder)):
    path = os.path.join(folder, file)
    print(f"\n--- {file} ---")
    if file.endswith('.docx'):
        try:
            doc = docx.Document(path)
            # print first 5 paragraphs
            count = 0
            for p in doc.paragraphs:
                if p.text.strip():
                    print(p.text.strip())
                    count += 1
                if count >= 10:
                    break
        except Exception as e:
            print("Error reading docx:", e)
    elif file.endswith('.pptx'):
        try:
            prs = Presentation(path)
            # print titles of first 5 slides
            count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        print(shape.text.strip().split('\n')[0]) # print first line of text
                        count += 1
                        break
                if count >= 8:
                    break
        except Exception as e:
            print("Error reading pptx:", e)
